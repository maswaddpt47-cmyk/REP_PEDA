"""
Génère les fiches mémo participant PPTX from scratch — format A4 recto-verso.
Structure fidèle au template Fiche_Memo.pptx : cartes numérotées, BSA, contact.
"""
import json, re, math
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT       = Path(__file__).parent.parent.parent
PPTX_DIR   = ROOT / "assets" / "pptx"
MEMOS_DIR  = ROOT / "assets" / "memos"
PLANS_FILE = ROOT / "assets" / "slide-plans.json"
BG_P1      = ROOT / "assets" / "bg-p1.png"   # fond page 1 (avec logo)
BG_P2      = ROOT / "assets" / "bg-p2.png"   # fond page 2+ (sans logo)

# ── Palette (template Fiche_Memo.pptx)
BLU     = RGBColor(0x43, 0x88, 0xBC)
TEA     = RGBColor(0x18, 0x7C, 0x88)
SKY     = RGBColor(0x5E, 0xB3, 0xD2)
DRK     = RGBColor(0x2B, 0x2B, 0x2B)
GRY     = RGBColor(0x6E, 0x6E, 0x6D)
WHT     = RGBColor(0xFF, 0xFF, 0xFF)
BG_BLU  = RGBColor(0xEC, 0xF5, 0xFF)
BG_TEA  = RGBColor(0xEB, 0xF7, 0xF6)
YEL     = RGBColor(0xF8, 0xA8, 0x24)
BRN     = RGBColor(0x5D, 0x40, 0x37)
BG_BSA  = RGBColor(0xFF, 0xFC, 0xE7)
BG_BDA  = RGBColor(0xEB, 0xF5, 0xFF)
BG_SUITE = RGBColor(0xEF, 0xF7, 0xFF)
SEP_CLR = RGBColor(0xDD, 0xDD, 0xDD)

STEP_PALETTE = [(BLU, BG_BLU), (TEA, BG_TEA)]

# ── Dimensions (cm, A4 portrait)
SW, SH   = 20.99, 29.70
ML       = 1.48
CW       = 18.03
ACCENT_W = 1.59
STEP_H   = 3.81
STEP_GAP = 0.18
SEP_X    = 3.35        # x-start of horizontal separator inside step
SEP_Y_IN = 1.17        # y within step where separator appears
SEP_W    = 15.98       # width of separator
HEADER_TOP  = 3.28
HEADER_H    = 1.42
FIRST_TOP   = 4.86
STEPS_PER_PAGE = 4
MAX_STEPS      = 8
CHARS_PER_LINE = 78   # estimation conservative à 10pt Arial sur 16.19 cm
LINE_H_CM      = 0.42 # hauteur d'une ligne wrappée à 10pt + espacement


def calc_step_h(bullets):
    """Hauteur dynamique d'un bloc step selon le nombre de lignes wrappées."""
    if not bullets:
        return STEP_H
    n_lines = sum(max(1, math.ceil(len(b) / CHARS_PER_LINE)) for b in bullets[:4])
    bullet_area = n_lines * LINE_H_CM + 0.35  # contenu + marge basse
    return max(SEP_Y_IN + 0.22 + bullet_area, STEP_H)


# ══════════════════════════════════════════════════════
# Primitives
# ══════════════════════════════════════════════════════

def add_bg_picture(slide, img_path, opacity=0.80):
    """Insère une image pleine page en fond avec transparence (opacity 0-1)."""
    from lxml import etree
    pic = slide.shapes.add_picture(str(img_path), Cm(0), Cm(0), Cm(SW), Cm(SH))
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    blip = pic._element.find(f'.//{{{a_ns}}}blip')
    if blip is not None:
        etree.SubElement(blip, f'{{{a_ns}}}alphaModFix', amt=str(int(opacity * 100000)))


def rect(slide, l, t, w, h, fill, line=None, lw=0.5):
    s = slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.text_frame.text = ""
    return s


def tbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top    = Cm(0.00)
    tf.margin_bottom = Cm(0.00)
    tf.margin_left   = Cm(0.05)
    tf.margin_right  = Cm(0.05)
    return tf


def p(tf, text, sz, bold=False, clr=DRK, align=PP_ALIGN.LEFT, first=False, italic=False):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = align
    para.space_before = Pt(0)
    para.space_after  = Pt(1)
    if text:
        run = para.add_run()
        run.text = text
        run.font.size      = Pt(sz)
        run.font.bold      = bold
        run.font.italic    = italic
        run.font.color.rgb = clr
        run.font.name      = "Arial"
    return para


# ══════════════════════════════════════════════════════
# Blocs de dessin
# ══════════════════════════════════════════════════════

def draw_memo_header(slide, titre, subtitle):
    """En-tête : titre + liste des thèmes."""
    tf = tbox(slide, ML, HEADER_TOP, CW, HEADER_H)
    titre_display = re.sub(r'^(Mon mémo[\s—\-]+)', '', titre)
    p(tf, f"Mon mémo — {titre_display}", 18, bold=True, clr=BLU, first=True)
    if subtitle:
        p(tf, subtitle, 10, clr=GRY)


def draw_memo_step(slide, num, titre, bullets, top, color_idx):
    """Carte étape : accent coloré + numéro + titre + bullets."""
    h  = calc_step_h(bullets)
    ac, bg = STEP_PALETTE[color_idx % 2]

    # Background + accent strip
    rect(slide, ML, top, CW, h, bg)
    rect(slide, ML, top, ACCENT_W, h, ac)

    # Horizontal separator (divides number/title from bullets)
    rect(slide, SEP_X, top + SEP_Y_IN, SEP_W, 0.02, SEP_CLR)

    # Step number — centred vertically in the title area
    num_h  = 0.80
    num_t  = top + (SEP_Y_IN - num_h) / 2
    tf_num = tbox(slide, ML, num_t, ACCENT_W, num_h)
    p(tf_num, str(num), 22, bold=True, clr=WHT, align=PP_ALIGN.CENTER, first=True)

    # Step title — right of accent, accent color
    tf_title = tbox(slide, ML + ACCENT_W + 0.20, top + 0.10, CW - ACCENT_W - 0.30, SEP_Y_IN - 0.14)
    p(tf_title, titre, 11, bold=True, clr=ac, first=True)

    # Bullets — below separator
    if bullets:
        tf_bull = tbox(slide, ML + ACCENT_W + 0.15, top + SEP_Y_IN + 0.12,
                       CW - ACCENT_W - 0.25, h - SEP_Y_IN - 0.22)
        first = True
        for b in bullets[:4]:
            p(tf_bull, f"→ {b}", 10, clr=DRK, first=first)
            first = False

    return top + h + STEP_GAP


def draw_suite_box(slide, page2_topics, top):
    """Encart 'Suite en page 2' affiché en bas de la page 1."""
    h = 1.83
    rect(slide, ML, top, CW, h, BG_SUITE, BLU, 0.5)
    tf = tbox(slide, ML + 2.80, top + 0.10, CW - 3.0, h - 0.15)
    topics_str = " | ".join(t for t in page2_topics[:3] if t)
    p(tf, f"→ Suite en page 2 : {topics_str}", 10.5, bold=True, clr=BLU, first=True)
    p(tf, "Votre Conseiller Numérique : 05 53 47 31 32  ·  conseiller-numerique@lotetgaronne.fr", 8.5, clr=GRY)


def draw_bsa(slide, titre, bullets, top):
    """Bloc 'Bon à savoir' (fond jaune, bordure ambre)."""
    n  = len(bullets)
    h  = 0.71 + 0.55 + n * 0.50 + 0.25
    h  = max(h, 2.20)

    rect(slide, ML, top, CW, h, BG_BSA, YEL, 0.7)
    # Separator line ambre au sein du bloc
    rect(slide, ML + 0.46, top + 0.71, CW - 0.92, 0.02, YEL)

    # Titre au-dessus du séparateur
    tf_t = tbox(slide, ML + 0.40, top + 0.08, CW - 0.55, 0.65)
    p(tf_t, titre, 12, bold=True, clr=BRN, first=True)

    # Bullets en dessous
    if bullets:
        tf_b = tbox(slide, ML + 0.40, top + 0.78, CW - 0.55, h - 0.88)
        first = True
        for b in bullets:
            p(tf_b, f"→ {b}", 10.5, clr=DRK, first=first)
            first = False

    return top + h + STEP_GAP


def draw_bda(slide, top):
    """Bloc 'Besoin d'aide ?' — contact CD47 (fond bleu clair)."""
    h = 4.20
    rect(slide, ML, top, CW, h, BG_BDA, BLU, 0.8)
    # Séparateur partiel (SKY)
    rect(slide, 6.17, top + 1.63, 8.64, 0.02, SKY)

    tf = tbox(slide, ML + 3.50, top + 0.18, 11.50, h - 0.30)
    p(tf, "Besoin d'aide ?", 13, bold=True, clr=BLU, first=True)
    p(tf, "Votre Conseiller Numérique est là pour vous accompagner !", 11, bold=True, clr=TEA)
    p(tf, "conseiller-numerique@lotetgaronne.fr  ·  Tél : 05 53 47 31 32", 11, clr=DRK)
    p(tf, "www.lotetgaronne.fr/inclusion-numerique", 10, clr=GRY)
    p(tf, "Direction des Systèmes d'Information — Service Inclusion Numérique — CD47", 9, clr=GRY)


def draw_footer(slide, page, total):
    """Ligne de pied de page version/pagination."""
    tf = tbox(slide, 6.62, 27.57, 7.77, 0.50)
    p(tf, f"v1.0 – 2025 – CD47 Service Inclusion Numérique  |  Page {page}/{total}",
      8, clr=GRY, align=PP_ALIGN.CENTER, first=True)


# ══════════════════════════════════════════════════════
# Extraction du contenu depuis le PPTX source
# ══════════════════════════════════════════════════════

def clean(text):
    return re.sub(r'[^\x00-\x7FÀ-ɏ''–—«»]', '', text).strip()

def trunc(t, n):
    return t[:n] + "…" if len(t) > n else t

def get_memo_bullets(slide, title, max_b=4):
    """Extrait les bullet points utiles d'une slide source."""
    bullets = []
    title_norm = re.sub(r'[^\w\s]', '', title).lower().strip()
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            t = clean(para.text.strip())
            if len(t) < 10: continue
            t_norm = re.sub(r'[^\w\s]', '', t).lower().strip()
            # Skip if essentially the slide title
            if t_norm == title_norm or t_norm in title_norm or title_norm in t_norm:
                continue
            bullets.append(t)
            if len(bullets) >= max_b: break
        if len(bullets) >= max_b: break
    return bullets


def extract_memo(pptx_path, slides_meta, version_indices=None):
    """Extrait contenu mémo depuis le PPTX + métadonnées."""
    prs = Presentation(pptx_path)

    idx_set = set(version_indices) if version_indices else None
    roles_steps = {'contenu', 'demo'}

    # Étapes : slides contenu non-optionnelles de la version
    steps = []
    for meta in slides_meta:
        if meta['role'] not in roles_steps: continue
        if meta.get('optionnel'): continue
        if idx_set and meta['n'] not in idx_set: continue
        slide = prs.slides[meta['n'] - 1]
        bullets = get_memo_bullets(slide, meta['titre'], 4)
        steps.append({
            'titre':   clean(meta['titre']),
            'bullets': bullets,
        })

    # Titre de l'atelier
    titre_raw = ''
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and len(t) < 120:
                titre_raw = t.split('\n')[0].strip(); break
    titre = trunc(clean(titre_raw), 60)

    # Blocs BSA : slides optionnelles absentes de la version
    bsa = []
    for meta in slides_meta:
        if not meta.get('optionnel'): continue
        if meta['role'] in ('titre', 'chapitre'): continue
        if idx_set and meta['n'] in idx_set: continue
        slide = prs.slides[meta['n'] - 1]
        bullets = get_memo_bullets(slide, meta['titre'], 4)
        bsa.append({
            'titre':   clean(meta['titre']),
            'bullets': bullets,
        })

    return {'titre': titre, 'steps': steps, 'bsa': bsa}


# ══════════════════════════════════════════════════════
# Construction du PPTX mémo
# ══════════════════════════════════════════════════════

def build_memo(memo_data, output_path):
    prs = Presentation()
    prs.slide_width  = Cm(SW)
    prs.slide_height = Cm(SH)

    _page = [0]
    def new_slide():
        _page[0] += 1
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHT
        bg = BG_P1 if _page[0] == 1 else BG_P2
        if bg.exists():
            add_bg_picture(s, bg)
        return s

    steps = memo_data['steps'][:MAX_STEPS]
    bsa   = memo_data['bsa']

    # Sous-titre : 3 premiers titres d'étapes
    subtitle = " | ".join(s['titre'] for s in steps[:3])

    # Découpage page 1 / page 2 basé sur la hauteur disponible
    # Page 1 : FIRST_TOP → ~25.5 cm (réserve suite box + bas de page)
    PAGE1_MAX = 25.50 - FIRST_TOP - 1.83 - 0.40  # suite_box_h + marge
    PAGE2_MAX = 25.50 - FIRST_TOP
    p1_steps, p2_steps = [], []
    used = 0.0
    for step in steps:
        h = calc_step_h(step['bullets'])
        if not p1_steps or used + h + STEP_GAP <= PAGE1_MAX:
            p1_steps.append(step)
            used += h + STEP_GAP
        else:
            p2_steps.append(step)
    needs_p2 = bool(p2_steps) or bool(bsa)
    total_pages = 2 if needs_p2 else 1

    # ─── PAGE 1 ──────────────────────────────────────
    s1 = new_slide()
    draw_memo_header(s1, memo_data['titre'], subtitle)

    top = FIRST_TOP
    for i, step in enumerate(p1_steps):
        top = draw_memo_step(s1, i + 1, step['titre'], step['bullets'], top, i)

    if needs_p2:
        # Suite box juste sous les étapes
        suite_topics = [s['titre'] for s in p2_steps[:3]]
        if bsa and len(suite_topics) < 3:
            suite_topics.append(bsa[0]['titre'])
        top = max(top + 0.20, FIRST_TOP + STEPS_PER_PAGE * (STEP_H + STEP_GAP) + 0.40)
        draw_suite_box(s1, suite_topics, top)
    else:
        # BSA + BDA sur page 1
        top += 0.30
        for b in bsa:
            top = draw_bsa(s1, b['titre'], b['bullets'], top)
        bda_top = max(top + 0.30, SH - 5.0)
        draw_bda(s1, bda_top)

    draw_footer(s1, 1, total_pages)

    # ─── PAGE 2 (si nécessaire) ───────────────────────
    if needs_p2:
        s2 = new_slide()
        draw_memo_header(s2, memo_data['titre'], subtitle)

        top = FIRST_TOP
        for i, step in enumerate(p2_steps):
            top = draw_memo_step(s2, len(p1_steps) + i + 1,
                                 step['titre'], step['bullets'], top,
                                 len(p1_steps) + i)

        top += 0.30
        for b in bsa:
            top = draw_bsa(s2, b['titre'], b['bullets'], top)

        bda_top = max(top + 0.30, SH - 5.0)
        draw_bda(s2, bda_top)
        draw_footer(s2, 2, total_pages)

    prs.save(output_path)


# ══════════════════════════════════════════════════════
# Point d'entrée
# ══════════════════════════════════════════════════════

def main():
    MEMOS_DIR.mkdir(exist_ok=True)
    with open(PLANS_FILE) as f:
        plans = json.load(f)

    generated = []
    for pptx_path in sorted(PPTX_DIR.glob("*.pptx")):
        key  = pptx_path.stem
        plan = plans.get(key, {})
        if not plan.get('analyse') or not plan.get('slides'):
            print(f"SKIP {key} — non analysé"); continue

        # Version source : 1h en priorité (étapes essentielles)
        vers   = plan.get('versions', {})
        v1h    = vers.get('1h', {}).get('indices')
        vfirst = vers.get(list(vers.keys())[0], {}).get('indices') if vers else None
        indices = v1h or vfirst

        try:
            memo_data = extract_memo(pptx_path, plan['slides'], indices)
            out = MEMOS_DIR / f"{key}-memo.pptx"
            build_memo(memo_data, out)
            n_steps = len(memo_data['steps'])
            n_bsa   = len(memo_data['bsa'])
            pages   = 2 if n_steps > STEPS_PER_PAGE or n_bsa else 1
            print(f"OK {key}: {n_steps} étapes | {n_bsa} BSA | {pages} page(s)")
            generated.append(key)
        except Exception as e:
            import traceback
            print(f"ERREUR {key}: {e}")
            traceback.print_exc()

    print(f"\n{len(generated)} mémo(s) généré(s) : {', '.join(generated)}")


if __name__ == "__main__":
    main()
