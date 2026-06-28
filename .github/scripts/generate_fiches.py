"""
Génère les fiches animateur PPTX from scratch — layout 100% dynamique.
Palette et structure fidèles au template original, sans dépendance au fichier .pptx.
N étapes déroulé, N lignes OPT (0-2), multi-pages si nécessaire.
"""
import json, re
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT       = Path(__file__).parent.parent.parent
PPTX_DIR   = ROOT / "assets" / "pptx"
FICHES_DIR = ROOT / "assets" / "fiches"
PLANS_FILE = ROOT / "assets" / "slide-plans.json"

# ── Palette (extraite du template original)
BLU    = RGBColor(0x43, 0x88, 0xBC)   # bleu objectifs / titres
TEA    = RGBColor(0x18, 0x7C, 0x88)   # teal matériel
SKY    = RGBColor(0x5E, 0xB3, 0xD2)   # bleu clair
ORG    = RGBColor(0xE6, 0x7D, 0x21)   # orange OPT
YEL    = RGBColor(0xF8, 0xA8, 0x24)   # jaune conseils
DRK    = RGBColor(0x2B, 0x2B, 0x2B)   # texte sombre
GRY    = RGBColor(0x6E, 0x6E, 0x6D)   # gris sous-titre
WHT    = RGBColor(0xFF, 0xFF, 0xFF)
BG_BLU = RGBColor(0xEC, 0xF5, 0xFF)   # fond objectifs
BG_TEA = RGBColor(0xEB, 0xF7, 0xF6)   # fond matériel
BG_STP = RGBColor(0xF8, 0xF9, 0xFD)   # fond étape déroulé
BD_STP = RGBColor(0xDD, 0xE6, 0xED)   # bordure étape
BG_OPT = RGBColor(0xFF, 0xF8, 0xEF)   # fond OPT
BG_CNS = RGBColor(0xFF, 0xFC, 0xE7)   # fond conseils

BADGE_CYCLE = [SKY, BLU, TEA, BLU, TEA, SKY]   # couleurs badges étapes (cycle)

# ── Dimensions (cm)
SW, SH   = 20.99, 29.70   # slide A4 portrait
ML       = 1.47            # marge gauche
CW       = 18.05           # largeur contenu

STEP_H      = 2.08    # hauteur d'une étape déroulé
OPT_HDR_H   = 0.55    # hauteur header section OPT
CONS_H      = 2.50    # hauteur section conseils


def opt_row_h(opt):
    """Hauteur dynamique d'une ligne OPT selon le nombre de bullets."""
    n = len(opt.get("bullets", []))
    return 1.05 + n * 0.42  # durée+titre+padding + n bullets

MATERIEL_DEFAULT = [
    "1 PC ou tablette par participant",
    "Connexion Wi-Fi stable",
    "Grand écran / vidéoprojecteur",
    "Fiches mémo (1 par participant)",
]
CONSEILS_DEFAULT = [
    "Adapter le rythme selon le niveau du groupe.",
    "Laisser les participants manipuler dès que possible.",
    "Prévoir un compte de démonstration pour ne pas exposer de données personnelles.",
    "Terminer en distribuant la fiche mémo — c'est le support qu'ils garderont.",
]


# ══════════════════════════════════════════════════════
# Primitives de dessin
# ══════════════════════════════════════════════════════

def rect(slide, l, t, w, h, fill, line=None, lw=0.5):
    """Rectangle coloré (en cm)."""
    s = slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.text_frame.text = ""
    return s


def tbox(slide, l, t, w, h):
    """Boîte de texte sans fond (en cm)."""
    tb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top    = Cm(0.00)
    tf.margin_bottom = Cm(0.00)
    tf.margin_left   = Cm(0.05)
    tf.margin_right  = Cm(0.05)
    return tf


def p(tf, text, sz, bold=False, clr=DRK, align=PP_ALIGN.LEFT, first=False, italic=False):
    """Ajoute un paragraphe+run dans un text_frame."""
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = align
    para.space_before = Pt(0)
    para.space_after  = Pt(1)
    if text:
        run = para.add_run()
        run.text = text
        run.font.size   = Pt(sz)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = clr
        run.font.name   = "Arial"
    return para


# ══════════════════════════════════════════════════════
# Blocs de contenu
# ══════════════════════════════════════════════════════

def draw_header(slide, titre, duree, top):
    tf = tbox(slide, ML, top, CW, 1.30)
    p(tf, titre, 17, bold=True, clr=BLU, align=PP_ALIGN.CENTER, first=True)
    p(tf, f"Fiche Action | {duree} | Guide d'animation — Conseillers Numériques CD47",
      8, clr=GRY, align=PP_ALIGN.CENTER)
    return top + 1.35


def draw_cols(slide, objectifs, materiel, top):
    col_w = 8.40
    gap   = 1.25
    col_h = 2.50

    # ── Objectifs
    rect(slide, ML, top, col_w, col_h, BG_BLU, BLU, 1.0)
    tf = tbox(slide, ML + 0.25, top + 0.15, col_w - 0.40, col_h - 0.20)
    p(tf, "Objectifs", 9.5, bold=True, clr=BLU, first=True)
    for obj in objectifs:
        p(tf, f"> {obj}" if obj else "", 8.5, clr=DRK)

    # ── Matériel
    mat_l = ML + col_w + gap
    rect(slide, mat_l, top, col_w, col_h, BG_TEA, TEA, 1.0)
    tf2 = tbox(slide, mat_l + 0.25, top + 0.15, col_w - 0.40, col_h - 0.20)
    p(tf2, "Matériel nécessaire", 9.5, bold=True, clr=TEA, first=True)
    for item in materiel:
        p(tf2, f"> {item}", 8.5, clr=DRK)

    return top + col_h + 0.20


def draw_step(slide, step, idx, top):
    clr = BADGE_CYCLE[idx % 6]

    # Fond de ligne + accent gauche coloré
    rect(slide, ML,        top, CW,   STEP_H, BG_STP, BD_STP, 0.5)
    rect(slide, ML,        top, 0.55, STEP_H, clr)

    # Badge heure (centré verticalement dans l'accent)
    bh = 0.44
    bt = top + (STEP_H - bh) / 2
    rect(slide, ML + 0.65, bt, 1.30, bh, clr)
    tb = tbox(slide, ML + 0.65, bt, 1.30, bh)
    p(tb, step["heure"], 8, bold=True, clr=WHT, align=PP_ALIGN.CENTER, first=True)

    # Contenu : durée + titre + bullet
    tf = tbox(slide, ML + 2.15, top + 0.18, CW - 2.35, STEP_H - 0.26)
    p(tf, f"{step['min']} min", 8, bold=True, clr=clr, first=True)
    p(tf, step["titre"], 9, bold=True, clr=DRK)
    if step.get("bullet"):
        p(tf, f". {step['bullet']}", 8, clr=DRK)

    return top + STEP_H


def draw_opt(slide, optionnels, top):
    """Dessine la section Blocs optionnels (0, 1 ou 2 lignes)."""
    if not optionnels:
        return top

    row_heights = [opt_row_h(o) for o in optionnels]
    total_h = OPT_HDR_H + sum(row_heights) + 0.10

    # Fond + bordure section
    rect(slide, ML, top, CW, total_h, BG_OPT, ORG, 1.5)

    # En-tête orange
    rect(slide, ML, top, CW, OPT_HDR_H, ORG)
    tf = tbox(slide, ML + 0.20, top + 0.08, CW - 0.40, OPT_HDR_H - 0.10)
    p(tf, "Blocs optionnels — À aborder si le temps le permet",
      8.5, bold=True, clr=WHT, align=PP_ALIGN.CENTER, first=True)

    row_top = top + OPT_HDR_H
    for i, opt in enumerate(optionnels):
        rh = row_heights[i]
        rect(slide, ML, row_top, CW, 0.04, ORG)   # séparateur fin

        tf2 = tbox(slide, ML + 0.30, row_top + 0.12, CW - 0.55, rh - 0.18)
        p(tf2, f"{opt['min']} min", 8, bold=True, clr=ORG, first=True)
        p(tf2, opt["titre"], 9, bold=True, clr=DRK)

        bullets = opt.get("bullets", [])
        if bullets:
            # Badge OPT. + premier bullet sur la même ligne
            para_obj = tf2.add_paragraph()
            para_obj.space_before = Pt(0)
            para_obj.space_after  = Pt(1)
            r0 = para_obj.add_run()
            r0.text = "OPT.  "
            r0.font.size = Pt(8); r0.font.bold = True
            r0.font.color.rgb = ORG; r0.font.name = "Arial"
            r1 = para_obj.add_run()
            r1.text = f". {bullets[0]}"
            r1.font.size = Pt(8); r1.font.color.rgb = DRK; r1.font.name = "Arial"
        for b in bullets[1:]:
            p(tf2, f". {b}", 8, clr=DRK)

        row_top += rh

    return top + total_h + 0.20


def draw_conseils(slide, conseils, top):
    rect(slide, ML, top, CW, 0.12, YEL)   # filet jaune
    tf = tbox(slide, ML + 0.25, top + 0.18, CW - 0.40, CONS_H - 0.20)
    p(tf, "Conseils pour une bonne animation", 9.5, bold=True, clr=DRK, first=True)
    for c in conseils:
        p(tf, f"> {c}", 8.5, clr=DRK)


# ══════════════════════════════════════════════════════
# Extraction des données depuis le PPTX source
# ══════════════════════════════════════════════════════

def clean(text):
    return re.sub(r'[^\x00-\x7FÀ-ɏ‘’–—«»]', '', text).strip()

def trunc(t, n):
    return t[:n] + "…" if len(t) > n else t

def get_slide_title(slide):
    if slide.shapes.title:
        t = slide.shapes.title.text.strip()
        if t: return t
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and len(t) < 120:
                return t.split("\n")[0].strip()
    return "(sans titre)"

def get_first_bullet(slide, title):
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t and t != title and len(t) > 10:
                t = clean(t)
                return t[:90] + ("…" if len(t) > 90 else "")
    return ""

def get_bullets(slide, title, max_b=3):
    bullets = []
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t and t != title and len(t) > 10:
                bullets.append(trunc(clean(t), 75))
                if len(bullets) >= max_b: return bullets
    return bullets


def extract_fiche(pptx_path, slides_meta, indices=None):
    """Extrait les données structurées depuis un PPTX + métadonnées."""
    prs = Presentation(pptx_path)
    titre = clean(get_slide_title(prs.slides[0]))

    idx_set = set(indices) if indices is not None else None
    pairs = [(s, m) for s, m in zip(prs.slides, slides_meta)
             if idx_set is None or m["n"] in idx_set]

    # Durée adaptée à la version
    duree_min = sum(m["min"] for _, m in pairs)
    h, mv = divmod(duree_min, 60)
    duree_str = (f"{h}h{mv:02d}" if h and mv else (f"{h}h" if h else f"{mv} min"))

    # Objectifs : 4 premières slides de contenu de la version
    content = [(s, m) for s, m in pairs
               if m["role"] in ("contenu", "demo", "chapitre") and m["n"] > 1]
    objectifs = [trunc(clean(m["titre"]), 55) for _, m in content[:4]]
    while len(objectifs) < 4: objectifs.append("")

    # Déroulé : slides non-optionnelles de la version, max 6 par page
    deroulé, cumtime = [], 0
    for slide, meta in pairs:
        if meta["role"] == "titre":
            cumtime += meta["min"]; continue
        if meta.get("optionnel"):
            cumtime += meta["min"]; continue
        h2, m2 = divmod(cumtime, 60)
        deroulé.append({
            "heure": f"{h2}:{m2:02d}",
            "min":   meta["min"],
            "titre": trunc(clean(meta["titre"]), 48),
            "bullet": get_first_bullet(slide, meta["titre"]),
        })
        cumtime += meta["min"]

    # Optionnels : slides optionnelles NON incluses dans cette version
    # (role "titre" exclu ; conclusion incluse si marquée optionnel)
    opt_raw = [(prs.slides[m["n"] - 1], m) for m in slides_meta
               if m.get("optionnel") and m["role"] != "titre"
               and (idx_set is None or m["n"] not in idx_set)]
    optionnels = [{"min": m["min"], "titre": trunc(clean(m["titre"]), 38),
                   "bullets": get_bullets(s, m["titre"], 3)}
                  for s, m in opt_raw]

    return {
        "titre": titre, "duree": duree_str,
        "objectifs": objectifs, "materiel": MATERIEL_DEFAULT,
        "deroulé": deroulé, "optionnels": optionnels,
        "conseils": CONSEILS_DEFAULT,
    }


# ══════════════════════════════════════════════════════
# Construction du PPTX
# ══════════════════════════════════════════════════════

def build_fiche(fiche_data, output_path):
    prs = Presentation()
    prs.slide_width  = Cm(SW)
    prs.slide_height = Cm(SH)

    def new_slide():
        s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHT
        return s

    slide = new_slide()
    top   = 3.40   # réserve l'espace logo/fond (à insérer plus tard)

    top = draw_header(slide, fiche_data["titre"], fiche_data["duree"], top)
    top += 0.15
    top = draw_cols(slide, fiche_data["objectifs"], fiche_data["materiel"], top)

    # Déroulé — saut de page si nécessaire
    for i, step in enumerate(fiche_data["deroulé"]):
        if top + STEP_H > SH - CONS_H - 0.60:
            slide = new_slide()
            top   = 1.50
        top = draw_step(slide, step, i, top)

    top += 0.25

    # OPT + conseils — saut de page si nécessaire
    n_opt   = len(fiche_data["optionnels"])
    opt_h   = (OPT_HDR_H + sum(opt_row_h(o) for o in fiche_data["optionnels"]) + 0.30) if n_opt else 0
    cons_h  = CONS_H + 0.30

    if top + opt_h + cons_h > SH - 0.30:
        slide = new_slide()
        top   = 1.50

    top = draw_opt(slide, fiche_data["optionnels"], top)

    # Conseils : ancrés en bas de la slide courante (pas d'espace vide)
    cons_top = max(top + 0.20, SH - CONS_H - 0.30)
    draw_conseils(slide, fiche_data["conseils"], cons_top)

    prs.save(output_path)


# ══════════════════════════════════════════════════════
# Point d'entrée
# ══════════════════════════════════════════════════════

def main():
    FICHES_DIR.mkdir(exist_ok=True)
    with open(PLANS_FILE) as f:
        plans = json.load(f)

    generated = []
    for pptx_path in sorted(PPTX_DIR.glob("*.pptx")):
        key  = pptx_path.stem
        plan = plans.get(key, {})
        if not plan.get("analyse") or not plan.get("slides"):
            print(f"SKIP {key} — non analysé"); continue

        versions = plan.get("versions") or {
            "complet": {"indices": [s["n"] for s in plan["slides"]]}
        }
        ok = []
        for vk, vdata in versions.items():
            try:
                fiche_data = extract_fiche(pptx_path, plan["slides"],
                                           indices=vdata["indices"])
                out = FICHES_DIR / f"{key}-fiche-{vk}.pptx"
                build_fiche(fiche_data, out)
                ok.append(f"{vk}({fiche_data['duree']})")
            except Exception as e:
                print(f"ERREUR {key}/{vk}: {e}")
        if ok:
            print(f"OK {key}: {', '.join(ok)}")
            generated.append(key)

    print(f"\n{len(generated)} atelier(s) traité(s) : {', '.join(generated)}")


if __name__ == "__main__":
    main()
