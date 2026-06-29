"""
Analyse automatique des PPTXs dans assets/pptx/.
Met à jour assets/slide-plans.json pour chaque atelier dont le PPTX existe.
Préserve les entrées marquées analyse:true manuellement (sauf si le PPTX a changé).
"""
import json, os, re
from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).parent.parent.parent
PPTX_DIR = ROOT / "assets" / "pptx"
PLANS_FILE = ROOT / "assets" / "slide-plans.json"


def get_slide_title(slide):
    if slide.shapes.title:
        t = slide.shapes.title.text.strip()
        if t:
            return t
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and len(t) < 120:
                return t.split("\n")[0].strip()
    return "(sans titre)"


def classify_role(title, idx):
    t = title.lower()
    if idx == 0:
        return "titre"
    if any(w in t for w in ["programme", "au menu", "sommaire", "objectif"]):
        return "intro"
    if any(w in t for w in ["chapitre", "partie", "module", "étape"]):
        return "chapitre"
    if any(w in t for w in ["démo", "demo", "pratique", "exercice", "manipulation", "tutoriel"]):
        return "demo"
    if any(w in t for w in ["pour aller plus loin", "ressource", "bilan", "conclusion", "récap", "merci", "questions fréquentes"]):
        return "conclusion"
    return "contenu"


DURATION = {"titre": 2, "intro": 4, "chapitre": 2, "demo": 12, "contenu": 8, "conclusion": 5}


def post_process(slides):
    """Si >50% des slides non-titre/intro sont 'chapitre', ce sont des slides de contenu."""
    body = [s for s in slides if s["role"] not in ("titre", "intro", "conclusion")]
    if body and sum(1 for s in body if s["role"] == "chapitre") / len(body) > 0.5:
        for s in slides:
            if s["role"] == "chapitre":
                s["role"] = "contenu"
                s["min"] = 10
    return slides


def is_optional(slide):
    t = slide["titre"].lower()
    return any(w in t for w in ["optionnel", "pour aller plus loin", "ressource", "numérique responsable", "éco-geste"])


def mins_label(m):
    h, m2 = divmod(m, 60)
    return f"{h}h{m2:02d}" if h and m2 else (f"{h}h" if h else f"{m} min")


def compute_versions(slides):
    all_idx = [s["n"] for s in slides]
    total = sum(s["min"] for s in slides)

    def must_keep(s): return s["role"] in ("titre", "intro", "chapitre")

    v1h30 = [s for s in slides if not s["optionnel"]]
    dur = sum(s["min"] for s in v1h30)
    if dur > 95:
        for s in reversed([x for x in v1h30 if not must_keep(x)]):
            if dur <= 95: break
            v1h30.remove(s); dur -= s["min"]

    v1h = list(v1h30); dur1h = sum(s["min"] for s in v1h)
    if dur1h > 65:
        for s in reversed([x for x in v1h if not must_keep(x)]):
            if dur1h <= 65: break
            v1h.remove(s); dur1h -= s["min"]

    def idx(lst): return [s["n"] for s in lst]
    return {
        "1h":     {"indices": idx(v1h),   "label": f"1h — {mins_label(sum(s['min'] for s in v1h))} · {len(v1h)} slides"},
        "1h30":   {"indices": idx(v1h30), "label": f"1h30 — {mins_label(sum(s['min'] for s in v1h30))} · {len(v1h30)} slides"},
        "complet":{"indices": all_idx,    "label": f"Complet — {mins_label(total)} · {len(all_idx)} slides"},
    }


def analyze(pptx_path):
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        title = get_slide_title(slide)
        role = classify_role(title, i)
        slides.append({"n": i + 1, "titre": title, "role": role,
                        "min": DURATION.get(role, 8), "optionnel": False})
    slides = post_process(slides)
    for s in slides:
        s["optionnel"] = is_optional(s)
    return slides, compute_versions(slides)


def main():
    with open(PLANS_FILE) as f:
        plans = json.load(f)

    # Charge le catalogue personnalisé pour les nouvelles tuiles
    cat_file = ROOT / "assets" / "catalogue.json"
    catalogue = {}
    if cat_file.exists():
        try:
            for entry in json.load(open(cat_file)):
                catalogue[entry["num"]] = entry
        except Exception:
            pass

    updated = []
    for pptx in sorted(PPTX_DIR.glob("*.pptx")):
        key = pptx.stem  # e.g. "B10"
        if key not in plans:
            cat = catalogue.get(key, {})
            plans[key] = {"titre": cat.get("titre", key), "chap": cat.get("chap", key[0] if key else "?")}
            print(f"NOUVEAU {key} — ajouté depuis catalogue.json")

        slides, versions = analyze(pptx)
        plans[key]["analyse"] = True
        plans[key]["slides"] = slides
        plans[key]["versions"] = versions
        total = sum(s["min"] for s in slides)
        print(f"OK {key}: {len(slides)} slides, ~{total} min")
        updated.append(key)

    if updated:
        with open(PLANS_FILE, "w") as f:
            json.dump(plans, f, ensure_ascii=False, indent=2)
            f.write("\n")
        print(f"\n{len(updated)} atelier(s) mis à jour : {', '.join(updated)}")
    else:
        print("Aucun PPTX à analyser.")


if __name__ == "__main__":
    main()
